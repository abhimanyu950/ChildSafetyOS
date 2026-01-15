"""
ChildSafetyOS Hackathon Presentation Generator
Generates a professional .pptx file using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb_color(r, g, b):
    """Helper to create RGB color tuple"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PRIMARY = rgb_color(102, 126, 234)      # Trust Blue #667EEA
    SECONDARY = rgb_color(118, 75, 162)     # Purple #764BA2
    SUCCESS = rgb_color(56, 161, 105)       # Green #38A169
    DANGER = rgb_color(229, 62, 62)         # Red #E53E3E
    DARK = rgb_color(26, 32, 44)            # Dark slate #1A202C
    WHITE = rgb_color(255, 255, 255)
    LIGHT_GRAY = rgb_color(226, 232, 240)
    
    def add_gradient_slide(title_text, subtitle_text=None):
        """Add a slide with gradient-like background"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background shape
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK
        bg.line.fill.background()
        
        # Title
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title_text
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = WHITE
        
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.5))
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle_text
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = rgb_color(160, 174, 192)
        
        return slide
    
    def add_bullet_points(slide, bullets, start_y=2.2, font_size=22):
        """Add bullet points to a slide"""
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(start_y), Inches(11.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.color.rgb = LIGHT_GRAY
            p.space_before = Pt(12)
            p.level = 0
    
    # ==================== SLIDE 1: Title ====================
    slide = add_gradient_slide("", "")
    
    # Shield emoji/icon area
    shield_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2), Inches(1))
    shield_tf = shield_box.text_frame
    shield_p = shield_tf.paragraphs[0]
    shield_p.text = "🛡️"
    shield_p.font.size = Pt(72)
    shield_p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "ChildSafetyOS"
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = "AI-Powered Parental Control at the OS Level"
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = rgb_color(160, 174, 192)
    sub_p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = '"Protecting children online, empowering parents everywhere."'
    tag_p.font.size = Pt(20)
    tag_p.font.italic = True
    tag_p.font.color.rgb = PRIMARY
    tag_p.alignment = PP_ALIGN.CENTER
    
    # ==================== SLIDE 2: Problem ====================
    slide = add_gradient_slide("The Digital World Isn't Built for Children", "The Problem")
    add_bullet_points(slide, [
        "93% of children aged 8-12 have access to smartphones",
        "56% have already encountered inappropriate content online",
        "Explicit content is just 2 clicks away on any search engine",
        "Traditional parental controls fail - kids bypass them in minutes",
        "Parents have zero visibility into what their children actually see"
    ])
    
    # ==================== SLIDE 3: Why Solutions Fail ====================
    slide = add_gradient_slide("Current Parental Controls Are Broken", "Why Existing Solutions Fail")
    add_bullet_points(slide, [
        "App-level blockers - Can be uninstalled or bypassed easily",
        "DNS filters only - Miss images, videos, embedded content",
        "Keyword blockers - No context understanding",
        "No real-time AI - Static lists can't catch new explicit content",
        "Zero visibility - Parents don't know what's actually happening"
    ])
    
    # ==================== SLIDE 4: Our Solution ====================
    slide = add_gradient_slide("ChildSafetyOS: Protection That Works", "Our Solution")
    add_bullet_points(slide, [
        "OS-Level Protection - Works at system level, cannot be bypassed",
        "Real-Time AI Analysis - Text and images analyzed instantly on-device",
        "Parent Dashboard - Live visibility into blocks, alerts, and activity",
        "Tamper-Proof - Settings and Play Store locked, only parents control"
    ])
    
    # ==================== SLIDE 5: Architecture ====================
    slide = add_gradient_slide("System Architecture", "How It Works")
    
    arch_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(4.5))
    arch_tf = arch_box.text_frame
    arch_tf.word_wrap = True
    
    arch_text = """CHILD'S DEVICE
Browser/App  -->  VPN Filter  -->  AI Engine  -->  Block/Allow

                        |
                        v
              
               CLOUD BACKEND
        Firebase - Events - Alerts

                        |
                        v

              PARENT DASHBOARD
          Stats - Logs - Controls"""
    arch_p = arch_tf.paragraphs[0]
    arch_p.text = arch_text
    arch_p.font.size = Pt(18)
    arch_p.font.name = "Consolas"
    arch_p.font.color.rgb = rgb_color(129, 230, 217)
    
    # ==================== SLIDE 6: Core Features ====================
    slide = add_gradient_slide("Core Features", "What ChildSafetyOS Delivers")
    add_bullet_points(slide, [
        "VPN Protection - All traffic filtered at OS level",
        "Domain Blocking - 100+ categories blocked (adult, gambling, drugs)",
        "AI Image Detection - NSFW images detected in real-time",
        "Text Analysis - Explicit keywords and context detected",
        "Settings Lock - Children cannot access Settings or Play Store",
        "Parent Dashboard - Live visibility and control",
        "Instant Alerts - Email notification for critical events"
    ], font_size=20)
    
    # ==================== SLIDE 7: How AI Is Used ====================
    slide = add_gradient_slide("How AI Is Used", "Smart Protection, Not Just Lists")
    add_bullet_points(slide, [
        "Text Analysis - Detects explicit keywords with context understanding",
        "Image Classification - TensorFlow Lite model runs ON DEVICE",
        "Privacy First - No images ever sent to cloud",
        "Age-Based Thresholds:",
        "   Child (Under 13): Maximum protection",
        "   Teen (13-17): Strong protection, allows mature themes",
        "   Adult (18+): Relaxed, still blocks explicit porn"
    ], font_size=20)
    
    # ==================== SLIDE 8: Data Flow ====================
    slide = add_gradient_slide("Live Data Flow", "What Happens in Milliseconds")
    add_bullet_points(slide, [
        "Step 1: Child opens a website",
        "Step 2: Traffic intercepted at OS level (VPN)",
        "Step 3: Domain checked against blocklist",
        "Step 4: AI analyzes images and text",
        "Step 5: Decision made in <500ms",
        "Step 6: ALLOW or BLOCK (with explanation)",
        "Step 7: Event logged - Parent dashboard updates",
        "Step 8: If critical - Parent gets email alert"
    ], font_size=20)
    
    # ==================== SLIDE 9: Tech Stack ====================
    slide = add_gradient_slide("Technology Stack", "Built for Production")
    add_bullet_points(slide, [
        "Android - Kotlin + Jetpack Compose (Native modern UI)",
        "VPN Service - Android VpnService API (OS-level interception)",
        "AI/ML - TensorFlow Lite (On-device classification)",
        "Backend - Firebase + Firestore (Real-time database)",
        "Cloud Functions - Node.js (Email alerts, server processing)",
        "Dashboard - HTML/JS + Firebase SDK (Parent monitoring)"
    ], font_size=20)
    
    # ==================== SLIDE 10: Privacy ====================
    slide = add_gradient_slide("Privacy & Safety by Design", "Protection Without Surveillance")
    add_bullet_points(slide, [
        "No images stored - AI runs on-device, images never uploaded",
        "Metadata only - We log 'blocked at URL' not actual content",
        "Encrypted communication - All sync uses HTTPS/TLS",
        "DPDP Compliant - Data deletion available on request",
        "No tracking - We don't sell or share any data",
        "Parent-controlled - Only parents decide what's blocked"
    ], font_size=20)
    
    # ==================== SLIDE 11: Dashboard ====================
    slide = add_gradient_slide("Parent Dashboard", "Real-Time Visibility & Control")
    add_bullet_points(slide, [
        "Live Statistics - Blocked vs Allowed, top domains, hourly activity",
        "Activity Logs - Every block recorded with reason and timestamp",
        "Controls - Age mode, category toggles, custom allowlist",
        "Alerts - Instant email for critical events, daily summary option",
        "",
        "This isn't surveillance - it's informed, modern parenting"
    ], font_size=22)
    
    # ==================== SLIDE 12: Technical Challenges ====================
    slide = add_gradient_slide("Technical Challenges Solved", "This Was Hard to Build")
    add_bullet_points(slide, [
        "OS-level interception - VPN Service API with packet routing",
        "Real-time ML on device - TensorFlow Lite optimized models",
        "Low latency decisions - <500ms including network + AI",
        "Accuracy vs speed - Perceptual hashing + ML pipeline",
        "Tamper resistance - Device Admin + Accessibility Service",
        "Battery efficiency - Smart caching, rate limiting"
    ], font_size=20)
    
    # ==================== SLIDE 13: Demo Status ====================
    slide = add_gradient_slide("Demo & Current Status", "What's Working Today")
    add_bullet_points(slide, [
        "VPN-based traffic filtering - Fully functional",
        "Domain blocking (100+ categories) - Complete",
        "AI image detection - On-device TFLite model working",
        "Safe Browser - Custom browser with content filtering",
        "Settings Lock - Accessibility Service blocking Settings/Play Store",
        "Firebase logging - Real-time event sync",
        "Parent dashboard - Live stats and activity logs",
        "Email alerts - Cloud Function triggers on critical events"
    ], font_size=18)
    
    # ==================== SLIDE 14: Impact ====================
    slide = add_gradient_slide("Impact & Use Cases", "Who Benefits?")
    add_bullet_points(slide, [
        "Parents - Peace of mind, visibility without surveillance",
        "Schools - Managed devices, compliance, central dashboard",
        "Device Manufacturers - Built-in safety for 'kid-safe' devices",
        "Government/Policy - Tool for child protection regulations"
    ])
    
    # ==================== SLIDE 15: Roadmap ====================
    slide = add_gradient_slide("Future Roadmap", "Where We're Going")
    add_bullet_points(slide, [
        "Q1 2026 - YouTube filtering, Screen time management",
        "Q2 2026 - Multi-device family support, iOS version",
        "Q3 2026 - School admin dashboard, Custom policy templates",
        "Beyond - Federated learning, Grooming detection, Geofencing"
    ])
    
    # ==================== SLIDE 16: Closing ====================
    slide = add_gradient_slide("", "")
    
    # Problem restate
    p1_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(0.8))
    p1_tf = p1_box.text_frame
    p1_p = p1_tf.paragraphs[0]
    p1_p.text = "The Problem: Children are exposed to harmful content. Current solutions don't work."
    p1_p.font.size = Pt(22)
    p1_p.font.color.rgb = rgb_color(252, 129, 129)
    
    # Solution restate
    p2_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(0.8))
    p2_tf = p2_box.text_frame
    p2_p = p2_tf.paragraphs[0]
    p2_p.text = "Our Solution: OS-level protection with real-time AI that cannot be bypassed."
    p2_p.font.size = Pt(22)
    p2_p.font.color.rgb = rgb_color(129, 230, 217)
    
    # Result restate
    p3_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.8))
    p3_tf = p3_box.text_frame
    p3_p = p3_tf.paragraphs[0]
    p3_p.text = "The Result: Parents get peace of mind. Children get safety. Privacy preserved."
    p3_p.font.size = Pt(22)
    p3_p.font.color.rgb = rgb_color(154, 230, 180)
    
    # Final tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.8))
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = '"Protecting the next generation, one device at a time."'
    tag_p.font.size = Pt(28)
    tag_p.font.italic = True
    tag_p.font.bold = True
    tag_p.font.color.rgb = PRIMARY
    tag_p.alignment = PP_ALIGN.CENTER
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    thanks_tf = thanks_box.text_frame
    thanks_p = thanks_tf.paragraphs[0]
    thanks_p.text = "Thank You"
    thanks_p.font.size = Pt(24)
    thanks_p.font.color.rgb = WHITE
    thanks_p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = r"c:\Users\rjii8\AndroidStudioProjects\ChildSafetyOS\ChildSafetyOS_Hackathon.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
